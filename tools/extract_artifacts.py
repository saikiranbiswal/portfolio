#!/usr/bin/env python3
"""Pre-convert portfolio artifacts into structured JSON for the in-site reader.

Run from the repo root:  python3 tools/extract_artifacts.py
Writes one JSON per artifact to assets/artifacts/content/<id>.json.
The reader (reader.html + js/reader.js) renders these in the site's design
system so visitors read artifacts without leaving the portfolio (no download
round-trip). Re-run whenever a source file in assets/artifacts/files changes.
"""
import csv
import json
import pathlib
import re
import sys

ROOT = pathlib.Path(__file__).resolve().parent.parent
OUT = ROOT / "assets" / "artifacts" / "content"
MAX_ROWS = 120          # cap per sheet so JSON stays light
MAX_COLS = 12


def cell_str(v):
    if v is None:
        return ""
    if isinstance(v, float) and v == int(v):
        return str(int(v))
    return str(v)


def extract_docx(path):
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(path)
    blocks = []

    def walk(parent):
        for child in parent.element.body.iterchildren():
            if child.tag.endswith("}p"):
                p = Paragraph(child, parent)
                text = p.text.strip()
                if not text:
                    continue
                style = (p.style.name if p.style is not None else "").lower()
                if "heading 1" in style or "title" in style:
                    kind = "h1"
                elif "heading 2" in style:
                    kind = "h2"
                elif "heading 3" in style or "heading 4" in style:
                    kind = "h3"
                elif "list" in style or p._p.xpath(".//w:numPr"):
                    kind = "li"
                else:
                    kind = "p"
                blocks.append({"t": kind, "x": text})
            elif child.tag.endswith("}tbl"):
                t = Table(child, parent)
                rows = []
                for r in t.rows[:MAX_ROWS]:
                    rows.append([c.text.strip() for c in r.cells[:MAX_COLS]])
                blocks.append({"t": "table", "rows": rows})

    walk(doc)
    return {"type": "doc", "blocks": blocks}


def extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows, truncated = [], False
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= MAX_ROWS:
                truncated = True
                break
            vals = [cell_str(v) for v in row[:MAX_COLS]]
            if any(v for v in vals):
                rows.append(vals)
        # trim fully-empty trailing columns
        width = max((len([v for v in r if v]) and max(j + 1 for j, v in enumerate(r) if v)) for r in rows) if rows else 0
        rows = [r[:width] for r in rows]
        if not rows:
            continue  # skip empty sheets — nothing to read
        sheets.append({"name": ws.title, "rows": rows, "truncated": truncated,
                       "total_rows": ws.max_row})
    return {"type": "sheet", "sheets": sheets}


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        title, body = "", []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                if not title and shape == slide.shapes.title:
                    title = text
                else:
                    body.append(text)
        if not title and body:
            title = body.pop(0)
        slides.append({"title": title, "body": body})
    return {"type": "deck", "slides": slides}


def extract_csv(path):
    rows = []
    with open(path, newline="", encoding="utf-8-sig") as f:
        for i, row in enumerate(csv.reader(f)):
            if i >= MAX_ROWS:
                break
            rows.append([c.strip() for c in row[:MAX_COLS]])
    return {"type": "sheet", "sheets": [{"name": pathlib.Path(path).stem,
                                         "rows": rows, "truncated": False,
                                         "total_rows": len(rows)}]}


def main():
    artifacts = json.loads((ROOT / "content" / "artifacts.json").read_text())["artifacts"]
    OUT.mkdir(parents=True, exist_ok=True)
    for a in artifacts:
        aid, rel = a["id"], a["file"]
        src = ROOT / rel
        ext = src.suffix.lower()
        if not src.exists():
            print(f"SKIP {aid}: missing {rel}", file=sys.stderr)
            continue
        if ext == ".docx":
            data = extract_docx(src)
        elif ext == ".xlsx":
            data = extract_xlsx(src)
        elif ext == ".pptx":
            data = extract_pptx(src)
        elif ext == ".csv":
            data = extract_csv(src)
        elif ext == ".svg":
            data = {"type": "svg", "src": rel}
        elif ext == ".pdf":
            data = {"type": "pdf", "src": rel}
        else:
            print(f"SKIP {aid}: unhandled {ext}", file=sys.stderr)
            continue
        data["id"] = aid
        out = OUT / f"{aid}.json"
        out.write_text(json.dumps(data, ensure_ascii=False))
        print(f"OK   {aid:38s} {data['type']:5s} -> {out.relative_to(ROOT)} ({out.stat().st_size//1024} KB)")


if __name__ == "__main__":
    main()
